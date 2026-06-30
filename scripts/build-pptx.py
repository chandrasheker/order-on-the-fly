#!/usr/bin/env python3
"""Build TableTap detailed owner presentation — comprehensive backup deck."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "presentation" / "screenshots"
OUT = ROOT / "presentation" / "TableTap-Restaurant-Owner-Detailed-Deck.pptx"

ORANGE = RGBColor(0xF9, 0x73, 0x16)
DARK = RGBColor(0x08, 0x08, 0x0E)
PANEL = RGBColor(0x16, 0x16, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA1, 0xA1, 0xAA)
EMERALD = RGBColor(0x34, 0xD3, 0x99)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)


def set_bg(slide, rgb=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12.1), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(8)


def add_bullets(slide, items, left=0.7, top=1.6, width=5.8, height=5.0, size=17):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)


def add_image(slide, name, left, top, width, height=None):
    path = SHOTS / f"{name}.png"
    if not path.exists():
        return False
    if height:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    return True


def add_footer(slide, text="TableTap · Complete Restaurant Platform"):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY


def card(slide, x, y, w, h, title, body):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = PANEL
    shp.line.color.rgb = RGBColor(55, 55, 72)
    tf = shp.text_frame
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if body:
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(6)


def build():
    SHOTS.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Title
    s = prs.slides.add_slide(blank)
    set_bg(s)
    t = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(3.0))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "TableTap"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p2 = tf.add_paragraph()
    p2.text = "Complete Restaurant Operating System"
    p2.font.size = Pt(28)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(12)
    p3 = tf.add_paragraph()
    p3.text = "QR dining · Kitchen · Payments · Takeaway · Swiggy/Zomato — one platform"
    p3.font.size = Pt(18)
    p3.font.color.rgb = GRAY
    p3.space_before = Pt(18)
    add_footer(s)

    # 2 Problem
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "The problem", "Restaurants juggle too many systems")
    add_bullets(s, [
        "Guests wait to order — slow during rush",
        "Kitchen runs on paper chits and shouting",
        "Swiggy, Zomato, and dine-in use different workflows",
        "Remote QR misuse wastes food and kitchen time",
        "Owners lack real-time revenue visibility",
    ])
    add_image(s, "10-staff-login", 7.0, 1.5, 4.8)
    add_footer(s)

    # 3 Solution
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "The TableTap solution", "One QR per table — full-service in one app")
    add_bullets(s, [
        "Guest scans QR → orders from phone (no app download)",
        "Kitchen sees all channels on one board (+ optional KDS)",
        "Staff dashboard with timers, alerts, payments",
        "Takeaway & delivery without a second POS",
        "Swiggy/Zomato orders flow in automatically",
        "Owner controls menu, QR, reports, integrations",
    ], width=6.0)
    add_image(s, "02-customer-menu", 7.0, 1.35, 4.6)
    add_footer(s)

    # 4 Customer journey
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Customer journey", "Secure check-in → order → pay → rewards")
    add_bullets(s, [
        "① Staff opens table when guests are seated",
        "② Guest scans QR and completes check-in",
        "③ Browse menu, add to cart, place order",
        "④ Track prep countdown on phone",
        "⑤ Pay via PhonePe QR or alert server",
        "⑥ Spin reward wheel while waiting",
    ], width=5.5)
    add_image(s, "03-customer-menu-scroll", 6.8, 1.2, 4.9)
    add_footer(s)

    # 5 Security
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Anti-misuse protection", "Saved links cannot order from home")
    add_bullets(s, [
        "Table must be opened by staff before ordering",
        "Rotating check-in + session limits per table",
        "Table closes when guests leave or bill is paid",
        "Prevents fake remote orders and kitchen waste",
    ], width=5.8)
    add_image(s, "01-customer-checkin", 7.0, 1.4, 4.5)
    add_footer(s)

    # 6 Staff dashboard
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Staff dashboard", "Live operations for every role")
    add_bullets(s, [
        "Active orders with prep timers & overdue alerts",
        "Sound alerts for new orders and help requests",
        "Mark preparing → ready → served",
        "Open/close table ordering with one tap",
        "Pending vs completed payments",
        "Reprint thermal receipts (premium)",
    ], width=5.5)
    add_image(s, "04-staff-dashboard", 6.5, 1.15, 6.0)
    add_footer(s)

    # 7 KDS
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Kitchen Display System", "Premium module — station-routed tickets")
    add_bullets(s, [
        "Routes items to Hot Kitchen, Grill, Bar, Cold",
        "Large-format UI designed for cooks",
        "Same stream: dine-in, takeaway, Swiggy, Zomato",
        "Optional auto-print kitchen chits (Bluetooth)",
        "Mark ready → pushes status to aggregators",
    ])
    card(s, 7.0, 1.5, 5.0, 1.2, "Enable: kds", "Super admin or CLI — no redeploy")
    card(s, 7.0, 3.0, 5.0, 1.2, "Route: /kitchen", "Cooks auto-redirect when enabled")
    add_footer(s)

    # 8 Floor + remote
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Floor plan & remote orders", "Premium modules")
    card(s, 0.8, 1.6, 5.5, 1.5, "Floor plan (floor_plan)", "Visual table map, server assignment, live bill, seat timers")
    card(s, 0.8, 3.4, 5.5, 1.5, "Remote orders (phone_orders)", "Walk-in, takeaway, delivery — one kitchen board")
    card(s, 0.8, 5.2, 5.5, 1.0, "Split bill (split_bill)", "Pay by item, split evenly, partial payments")
    add_image(s, "05-table-ordering-panel", 6.8, 1.5, 5.5, 2.8)
    add_footer(s)

    # 9 Payments
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Payments", "PhonePe QR + staff verification + thermal receipts")
    add_bullets(s, [
        "Upload PhonePe static QR once in admin",
        "Guest taps Pay → QR on their phone",
        "Staff marks paid → table unlocks",
        "Revenue reports count only confirmed payments",
        "Bluetooth ESC/POS receipt on full pay (premium)",
        "GSTIN on receipts when gst_receipts enabled",
    ], width=5.8)
    img = "06-pending-payments" if (SHOTS / "06-pending-payments.png").exists() else "04-staff-dashboard"
    add_image(s, img, 6.8, 1.2, 5.8)
    add_footer(s)

    # 10 Swiggy Zomato
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Swiggy & Zomato integration", "Premium aggregator_inbox — automatic ready mode")
    add_bullets(s, [
        "Owner saves outlet ID + API key in Admin → Integrations",
        "TableTap generates webhook URL + secret",
        "Partner team activates webhook (one-time)",
        "Orders appear on kitchen board — no manual entry",
        "Menu syncs outbound when prices/stock change",
        "Ready / picked up / delivered pushed back to platform",
    ], width=6.2)
    card(s, 7.2, 1.5, 5.0, 2.0, "Full loop", "Ingest → confirm → prepare → ready → picked up. See AGGREGATOR_SETUP.md")
    card(s, 7.2, 3.8, 5.0, 1.2, "Menu mapping", "Set swiggyItemId / zomatoItemId or match names")
    add_footer(s)

    # 11 Admin
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Owner admin tools", "Menu · QR · Reports · Integrations")
    add_image(s, "07-admin-qr-codes", 0.8, 1.25, 5.8)
    add_image(s, "08-admin-menu", 6.85, 1.25, 5.8)
    add_footer(s)

    # 12 Reports
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Reports & team performance", "Daily CSV + accountability")
    add_image(s, "09-admin-reports", 0.8, 1.25, 11.8)
    add_footer(s)

    # 13 Core vs Premium
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Core vs Premium", "Flexible tiers — toggle live without downtime")
    add_bullets(s, [
        "CORE: QR ordering, staff dashboard, menu, payments, PhonePe QR, sessions, rewards, reports",
        "PREMIUM: KDS, floor plan, split bill, thermal receipts, phone orders, GST, aggregator sync, staff performance",
        "Super admin at /platform/login — hidden from restaurant owners",
        "Enable via CLI: npx tsx scripts/enable-premium-features.ts --slug SLUG --all",
    ], width=11.5, top=1.5)
    add_footer(s)

    # 14 Roles
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Role-based permissions", "Owner · Manager · Cook · Server · Platform admin")
    add_bullets(s, [
        "Owner / Manager — dashboard, menu, QR, reports, integrations",
        "Cook — KDS, prepare items, mark unavailable",
        "Server — floor plan, serve, open/close tables, payments",
        "Platform admin — staff slots, premium toggles, credential export",
    ])
    add_footer(s)

    # 15 Setup
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Technical setup", "See SETUP_GUIDE.md for complete instructions")
    add_bullets(s, [
        "1. npm install && npm run setup -- --start",
        "2. Edit restaurant.config.json (name, staff, menu, tables)",
        "3. Set NEXT_PUBLIC_APP_URL to your LAN IP or domain",
        "4. Print QR codes from admin",
        "5. Enable premium modules from /platform/login",
        "6. Docker: docker compose up --build (PostgreSQL production)",
        "7. Swiggy/Zomato: Admin → Integrations + partner webhook registration",
    ], width=11.0)
    add_footer(s)

    # 16 Benefits
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Benefits for restaurant owners", "Why TableTap pays for itself")
    items = [
        ("Faster service", "Guests order without waiting for staff"),
        ("One system", "Dine-in + takeaway + Swiggy + Zomato"),
        ("Less waste", "Stock control + remote order blocking"),
        ("Higher turnover", "Quicker payments, faster table reset"),
        ("Kitchen clarity", "KDS + timers replace paper chaos"),
        ("Full control", "You own menu, staff, QR, reports, integrations"),
    ]
    y = 1.5
    for title, desc in items:
        box = s.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.5), Inches(0.7))
        tf = box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"✓  {title} — "
        run.font.bold = True
        run.font.size = Pt(18)
        run.font.color.rgb = EMERALD
        run2 = p.add_run()
        run2.text = desc
        run2.font.size = Pt(18)
        run2.font.color.rgb = WHITE
        y += 0.75
    add_footer(s)

    # 17 Demo
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "Live demo credentials", "Varanasi Restaurant demo")
    add_bullets(s, [
        "Staff login: owner@varanasi.com / admin123",
        "Super admin: admin@varanasi.com / admin@varanasi → /platform/login",
        "Customer: scan Table 1 QR from staff dashboard",
        "",
        "Documentation:",
        "  SETUP_GUIDE.md — complete setup",
        "  AGGREGATOR_SETUP.md — Swiggy/Zomato",
        "  PREMIUM_FEATURES.md — module list",
    ], width=11.0)
    add_footer(s, "Let's run a live demo at your restaurant.")

    # 18 Thank you
    s = prs.slides.add_slide(blank)
    set_bg(s)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), Inches(13.333), Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(2.5))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank you"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "TableTap — Scan · Order · Serve · Grow"
    p2.font.size = Pt(24)
    p2.font.color.rgb = ORANGE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)
    p3 = tf.add_paragraph()
    p3.text = "Complete setup guide included · Ready for your restaurant"
    p3.font.size = Pt(16)
    p3.font.color.rgb = GRAY
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(12)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved {OUT}")


if __name__ == "__main__":
    build()
