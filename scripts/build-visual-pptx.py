#!/usr/bin/env python3
"""Build a premium visual-first TableTap pitch deck — works with or without screenshots."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "presentation" / "screenshots"
OUT = ROOT / "presentation" / "TableTap-Restaurant-Owner-Visual-Deck.pptx"
DETAILED_OUT = ROOT / "presentation" / "TableTap-Restaurant-Owner-Detailed-Deck.pptx"
LEGACY_OUT = ROOT / "presentation" / "TableTap-Restaurant-Owner-Deck.pptx"

# Palette — dark premium restaurant tech
DARK = RGBColor(8, 8, 14)
PANEL = RGBColor(22, 22, 34)
PANEL_2 = RGBColor(32, 32, 48)
PANEL_3 = RGBColor(42, 38, 58)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(160, 160, 175)
ORANGE = RGBColor(249, 115, 22)
AMBER = RGBColor(245, 158, 11)
GREEN = RGBColor(52, 211, 153)
RED = RGBColor(248, 113, 113)
PURPLE = RGBColor(168, 85, 247)
CYAN = RGBColor(34, 211, 238)
BLUE = RGBColor(96, 165, 250)
BORDER = RGBColor(55, 55, 72)


def bg(slide, accent=PURPLE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK
    for x, y, size, color, trans in [
        (-1.0, -0.8, 3.4, ORANGE, 78),
        (10.5, -0.2, 2.8, accent, 74),
        (9.2, 5.5, 2.6, GREEN, 76),
        (-0.5, 5.8, 2.2, CYAN, 80),
    ]:
        blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        blob.fill.solid()
        blob.fill.fore_color.rgb = color
        blob.fill.transparency = trans
        blob.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def textbox(slide, text, x, y, w, h, size=24, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def title(slide, text, kicker=None, subtitle=None):
    if kicker:
        textbox(slide, kicker.upper(), 0.65, 0.42, 6.0, 0.35, 10, ORANGE, True)
    textbox(slide, text, 0.62, 0.78, 8.5, 0.85, 36, WHITE, True)
    if subtitle:
        textbox(slide, subtitle, 0.65, 1.55, 8.0, 0.55, 16, MUTED)


def pill(slide, text, x, y, w, color=ORANGE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.42))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    p = shp.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return shp


def card(slide, x, y, w, h, heading, caption="", color=PANEL, icon=""):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = BORDER
    tf = shp.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.14)
    p = tf.paragraphs[0]
    p.text = f"{icon}  {heading}" if icon else heading
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if caption:
        p2 = tf.add_paragraph()
        p2.text = caption
        p2.font.size = Pt(11)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(8)
    return shp


def screenshot(slide, name, x, y, w, h=None, frame=True):
    path = SHOTS / f"{name}.png"
    if not path.exists():
        return None
    if frame:
        pad = 0.1
        fh = h if h else w * 0.62
        fr = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x - pad),
            Inches(y - pad),
            Inches(w + pad * 2),
            Inches(fh + pad * 2),
        )
        fr.fill.solid()
        fr.fill.fore_color.rgb = RGBColor(4, 4, 8)
        fr.line.color.rgb = BORDER
    if h:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


def phone(slide, name, x, y, w):
    frame_h = w * 2.05
    fr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(frame_h))
    fr.fill.solid()
    fr.fill.fore_color.rgb = RGBColor(2, 2, 6)
    fr.line.color.rgb = RGBColor(80, 80, 96)
    pic = screenshot(slide, name, x + 0.12, y + 0.12, w - 0.24, frame_h - 0.24, frame=False)
    if not pic:
        textbox(slide, "Guest\nmenu", x + 0.2, y + frame_h * 0.35, w - 0.4, 1.0, 14, MUTED, True, PP_ALIGN.CENTER)
    return fr


def footer(slide, n, total=18):
    textbox(slide, f"TableTap  ·  {n}/{total}", 0.65, 7.05, 4.0, 0.25, 9, MUTED)
    textbox(slide, f"{n:02d}", 12.35, 6.92, 0.5, 0.25, 9, MUTED, True, PP_ALIGN.RIGHT)


def stat_row(slide, stats, y=5.2):
    w = 2.05
    gap = 0.35
    for i, (val, label, color) in enumerate(stats):
        x = 0.75 + i * (w + gap)
        card(slide, x, y, w, 1.05, val, label, PANEL_2)
        # recolor top stat
        shp = slide.shapes[-1]
        shp.text_frame.paragraphs[0].font.color.rgb = color


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    TOTAL = 18

    # ── 1 Hero ──
    s = prs.slides.add_slide(blank)
    bg(s, PURPLE)
    pill(s, "RESTAURANT OPERATING SYSTEM", 0.72, 0.85, 2.55)
    textbox(s, "TableTap", 0.65, 1.65, 6.0, 1.0, 58, WHITE, True)
    textbox(s, "Scan · Order · Serve · Grow", 0.68, 2.75, 5.8, 0.55, 26, ORANGE, True)
    textbox(
        s,
        "QR dining + kitchen ops + payments + Swiggy/Zomato — one platform, zero app download for guests.",
        0.7,
        3.45,
        5.2,
        1.0,
        16,
        MUTED,
    )
    card(s, 0.75, 4.85, 1.55, 0.95, "No guest app", "Scan & order", PANEL_2, "📱")
    card(s, 2.45, 4.85, 1.55, 0.95, "Live kitchen", "Timers + KDS", PANEL_2, "👨‍🍳")
    card(s, 4.15, 4.85, 1.55, 0.95, "Aggregator sync", "Swiggy · Zomato", PANEL_2, "🛵")
    card(s, 5.85, 4.85, 1.55, 0.95, "Owner control", "Menu · QR · Reports", PANEL_2, "📊")
    phone(s, "02-customer-menu", 8.35, 0.65, 3.05)
    footer(s, 1, TOTAL)

    # ── 2 Pain points ──
    s = prs.slides.add_slide(blank)
    bg(s, RED)
    title(s, "Problems every restaurant faces", "owner pain points", "Before TableTap")
    pains = [
        ("Slow ordering", "Guests wait for staff during rush"),
        ("Kitchen chaos", "Paper chits, shouting, mistakes"),
        ("Payment delays", "Tables blocked, turnover drops"),
        ("Remote fraud", "Saved QR links order from home"),
        ("Aggregator overload", "Swiggy + Zomato + dine-in = 3 systems"),
        ("No visibility", "Owner finds out revenue at day end"),
    ]
    for i, (h, c) in enumerate(pains):
        x = 0.75 + (i % 3) * 4.05
        y = 2.0 + (i // 3) * 1.55
        card(s, x, y, 3.55, 1.25, h, c, PANEL_2)
    if screenshot(s, "04-staff-dashboard", 8.8, 2.0, 3.8, 2.5):
        pass
    footer(s, 2, TOTAL)

    # ── 3 Platform overview ──
    s = prs.slides.add_slide(blank)
    bg(s, CYAN)
    title(s, "One platform — full service", "complete solution")
    modules = [
        ("Guest QR ordering", "Secure check-in, cart, rewards", GREEN),
        ("Staff dashboard", "Live board, alerts, payments", ORANGE),
        ("Kitchen KDS", "Station-routed tickets", RED),
        ("Floor plan", "Table map + server assign", BLUE),
        ("Takeaway & delivery", "Walk-in, pack, ship", PURPLE),
        ("Swiggy / Zomato", "Auto ingest + menu sync", AMBER),
    ]
    for i, (h, c, col) in enumerate(modules):
        x = 0.75 + (i % 3) * 4.05
        y = 2.05 + (i // 3) * 1.75
        card(s, x, y, 3.55, 1.35, h, c, PANEL_3)
    textbox(s, "Core included · Premium modules toggle on without downtime", 1.5, 5.85, 10.0, 0.4, 18, GREEN, True, PP_ALIGN.CENTER)
    footer(s, 3, TOTAL)

    # ── 4 Customer experience ──
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "Guest experience guests love", "no app · no login")
    phone(s, "01-customer-checkin", 0.75, 1.35, 2.65)
    phone(s, "02-customer-menu", 3.65, 1.35, 2.65)
    phone(s, "03-customer-menu-scroll", 6.55, 1.35, 2.65)
    card(s, 9.55, 1.65, 2.95, 0.85, "Secure check-in", "Staff opens table first")
    card(s, 9.55, 2.75, 2.95, 0.85, "Browse & order", "Categories, prep times")
    card(s, 9.55, 3.85, 2.95, 0.85, "Track & pay", "Timer + PhonePe QR")
    card(s, 9.55, 4.95, 2.95, 0.85, "Rewards", "Spin wheel while waiting")
    footer(s, 4, TOTAL)

    # ── 5 Security ──
    s = prs.slides.add_slide(blank)
    bg(s, GREEN)
    title(s, "Stops remote order misuse", "security built in")
    steps = [
        ("1", "Staff opens table", "Only seated guests can order"),
        ("2", "Guest scans QR", "Check-in cookie issued"),
        ("3", "Orders flow", "Saved links stay blocked"),
        ("4", "Table closes", "Auto after payment or manual"),
    ]
    for i, (num, h, c) in enumerate(steps):
        x = 0.75 + i * 3.05
        textbox(s, num, x + 0.15, 2.15, 0.5, 0.5, 28, ORANGE, True)
        card(s, x, 2.05, 2.65, 1.15, h, c, PANEL_2)
    pill(s, "LESS FOOD WASTE", 1.0, 4.0, 2.0, GREEN)
    pill(s, "STAFF IN CONTROL", 3.5, 4.0, 2.1, AMBER)
    pill(s, "SESSION LIMITS", 6.0, 4.0, 1.9, CYAN)
    screenshot(s, "01-customer-checkin", 4.5, 4.35, 3.5, 2.2)
    footer(s, 5, TOTAL)

    # ── 6 Staff dashboard ──
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "Staff dashboard — mission control", "live operations")
    screenshot(s, "04-staff-dashboard", 0.7, 1.5, 7.35, 4.35)
    card(s, 8.45, 1.65, 3.85, 0.88, "Live order board", "All tables, all channels")
    card(s, 8.45, 2.78, 3.85, 0.88, "Prep timers", "Overdue alerts + sound")
    card(s, 8.45, 3.91, 3.85, 0.88, "Mark ready / served", "Full item lifecycle")
    card(s, 8.45, 5.04, 3.85, 0.88, "Open / close tables", "One tap control")
    footer(s, 6, TOTAL)

    # ── 7 KDS ──
    s = prs.slides.add_slide(blank)
    bg(s, RED)
    title(s, "Kitchen Display System", "premium · ends paper chits")
    card(s, 0.75, 2.0, 3.6, 1.4, "Station routing", "Hot · Grill · Bar · Cold", PANEL_2, "🔥")
    card(s, 4.65, 2.0, 3.6, 1.4, "Large ticket UI", "Cook-friendly, high contrast", PANEL_2, "📺")
    card(s, 8.55, 2.0, 3.6, 1.4, "Same order stream", "Dine-in + takeaway + Swiggy", PANEL_2, "🔄")
    card(s, 0.75, 3.75, 5.5, 1.4, "Auto kitchen chit print", "Bluetooth ESC/POS when enabled", PANEL_2, "🖨️")
    card(s, 6.55, 3.75, 5.5, 1.4, "Status → aggregators", "Ready / picked up pushed to Swiggy/Zomato", PANEL_2, "✅")
    pill(s, "PREMIUM MODULE: kds", 0.85, 5.55, 2.4, ORANGE)
    footer(s, 7, TOTAL)

    # ── 8 Floor plan ──
    s = prs.slides.add_slide(blank)
    bg(s, BLUE)
    title(s, "Floor plan for servers", "premium · visual table map")
    screenshot(s, "05-table-ordering-panel", 0.75, 1.45, 11.85, 3.05)
    card(s, 0.95, 5.0, 2.85, 0.9, "Seat & assign", "Server ownership")
    card(s, 4.05, 5.0, 2.85, 0.9, "Live bill total", "Per-table revenue")
    card(s, 7.15, 5.0, 2.85, 0.9, "Timer per table", "Turnover insight")
    card(s, 10.25, 5.0, 2.35, 0.9, "Open / close", "Table control")
    footer(s, 8, TOTAL)

    # ── 9 Payments ──
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "Payments & split bill", "cash · UPI · partial pay")
    img = screenshot(s, "06-pending-payments", 0.75, 1.35, 6.35, 3.85)
    if not img:
        screenshot(s, "04-staff-dashboard", 0.75, 1.35, 6.35, 3.85)
    card(s, 7.55, 1.55, 4.5, 0.9, "PhonePe QR on guest phone", "Upload once in admin")
    card(s, 7.55, 2.7, 4.5, 0.9, "Split by item or evenly", "Premium split_bill")
    card(s, 7.55, 3.85, 4.5, 0.9, "Pending vs completed", "Revenue = paid only")
    card(s, 7.55, 5.0, 4.5, 0.9, "Thermal receipt on pay", "Bluetooth ESC/POS")
    footer(s, 9, TOTAL)

    # ── 10 Remote orders ──
    s = prs.slides.add_slide(blank)
    bg(s, PURPLE)
    title(s, "Takeaway & delivery", "premium · no second POS")
    modes = [
        ("Walk-in / table", "Staff places order for seated guest"),
        ("Takeaway", "Pack and hand over — no table"),
        ("Delivery", "Phone + address notes"),
    ]
    for i, (h, c) in enumerate(modes):
        card(s, 0.85, 1.85 + i * 1.45, 4.8, 1.15, h, c, PANEL_2)
    card(s, 6.2, 1.85, 6.0, 3.55, "One kitchen board", "Dine-in + takeaway + delivery + aggregators — same workflow, same printers, same reports.", PANEL_3)
    pill(s, "PREMIUM: phone_orders", 6.35, 5.65, 2.5, PURPLE)
    footer(s, 10, TOTAL)

    # ── 11 Swiggy / Zomato ──
    s = prs.slides.add_slide(blank)
    bg(s, AMBER)
    title(s, "Swiggy & Zomato — automatic", "premium · ready mode")
    flow = [
        ("Save credentials", "Admin → Integrations"),
        ("Register webhook", "One-time with partner team"),
        ("Orders auto-flow", "Kitchen board — zero manual entry"),
        ("Menu sync out", "Price & stock pushed back"),
        ("Status sync back", "Ready → picked up → delivered"),
    ]
    for i, (h, c) in enumerate(flow):
        y = 1.75 + i * 0.95
        textbox(s, str(i + 1), 0.85, y + 0.08, 0.4, 0.4, 18, ORANGE, True)
        card(s, 1.35, y, 4.5, 0.75, h, c, PANEL_2)
    card(s, 6.3, 1.75, 6.0, 3.8, "Full loop closed", "Ingest → confirm → prepare → ready → picked up. Owner never re-enters orders in Swiggy/Zomato apps.", PANEL_3, "🛵")
    pill(s, "PREMIUM: aggregator_inbox", 6.45, 5.75, 2.8, AMBER)
    footer(s, 11, TOTAL)

    # ── 11 Admin tools ──
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "Owner admin tools", "menu · QR · reports · integrations")
    screenshot(s, "07-admin-qr-codes", 0.7, 1.25, 5.85, 3.55)
    screenshot(s, "08-admin-menu", 6.85, 1.25, 5.85, 3.55)
    pill(s, "PRINT TABLE QR", 1.1, 5.15, 1.95, ORANGE)
    pill(s, "PHONEPE QR UPLOAD", 3.35, 5.15, 2.35, GREEN)
    pill(s, "EDIT MENU & STOCK", 6.95, 5.15, 2.35, CYAN)
    pill(s, "SWIGGY / ZOMATO", 9.65, 5.15, 2.35, AMBER)
    footer(s, 12, TOTAL)

    # ── 13 Reports ──
    s = prs.slides.add_slide(blank)
    bg(s, GREEN)
    title(s, "Reports & team performance", "data owners actually use")
    screenshot(s, "09-admin-reports", 0.75, 1.3, 6.85, 4.25)
    card(s, 8.15, 1.55, 4.2, 0.85, "Daily revenue & orders", "CSV export")
    card(s, 8.15, 2.65, 4.2, 0.85, "Item & table breakdown", "What sells")
    card(s, 8.15, 3.75, 4.2, 0.85, "Team performance", "Prep · serve · collect")
    card(s, 8.15, 4.85, 4.2, 0.85, "Aggregator last sync", "Menu + order health")
    footer(s, 13, TOTAL)

    # ── 14 Core vs Premium ──
    s = prs.slides.add_slide(blank)
    bg(s, PURPLE)
    title(s, "Flexible pricing tiers", "core included · premium add-ons")
    textbox(s, "CORE (always on)", 0.85, 1.75, 5.5, 0.35, 14, GREEN, True)
    core = "QR ordering · Staff dashboard · Menu admin · Table QR · Payments · PhonePe QR · Sessions · Rewards · Reports · Alerts"
    textbox(s, core, 0.85, 2.15, 5.8, 1.5, 12, MUTED)
    textbox(s, "PREMIUM (toggle per restaurant)", 7.0, 1.75, 5.5, 0.35, 14, ORANGE, True)
    prem = "KDS · Floor plan · Split bill · Thermal receipts · Takeaway/delivery · Swiggy/Zomato sync · GST receipts · Staff performance"
    textbox(s, prem, 7.0, 2.15, 5.8, 1.5, 12, MUTED)
    card(s, 0.85, 4.15, 3.6, 1.2, "Starter", "All core features", PANEL_2)
    card(s, 4.85, 4.15, 3.6, 1.2, "Pro", "Core + chosen premium", PANEL_3)
    card(s, 8.85, 4.15, 3.6, 1.2, "Enterprise", "Full platform + roadmap", PANEL_2)
    textbox(s, "Toggle modules live — no downtime, no redeploy", 2.5, 5.75, 8.5, 0.4, 16, CYAN, True, PP_ALIGN.CENTER)
    footer(s, 14, TOTAL)

    # ── 15 Roles ──
    s = prs.slides.add_slide(blank)
    bg(s)
    title(s, "Right access for every role", "permissions")
    roles = [
        ("Owner", "Reports · admin · integrations", "👑"),
        ("Manager", "Menu · payments · QR", "📋"),
        ("Cook", "KDS · prepare · stock", "👨‍🍳"),
        ("Server", "Floor · serve · tables", "🍽️"),
    ]
    for i, (h, c, icon) in enumerate(roles):
        card(s, 0.85 + i * 3.05, 2.15, 2.65, 1.55, h, c, PANEL_2, icon)
    textbox(s, "Super admin (hidden) toggles premium per restaurant", 2.2, 4.5, 8.8, 0.4, 18, PURPLE, True, PP_ALIGN.CENTER)
    footer(s, 15, TOTAL)

    # ── 16 Setup ──
    s = prs.slides.add_slide(blank)
    bg(s, CYAN)
    title(s, "Go live in one day", "setup")
    steps = [
        ("① Config", "One JSON file — name, staff, menu, tables"),
        ("② Deploy", "npm run setup or Docker Compose"),
        ("③ Print QR", "Place on every table"),
        ("④ Train staff", "Open table → scan → serve → paid"),
        ("⑤ Aggregators", "Credentials + webhook (optional)"),
    ]
    for i, (h, c) in enumerate(steps):
        card(s, 0.75 + (i % 3) * 4.05, 2.05 + (i // 3) * 1.65, 3.55, 1.25, h, c, PANEL_2)
    textbox(s, "Full guide: SETUP_GUIDE.md in the repo", 2.5, 5.65, 8.5, 0.4, 15, MUTED, False, PP_ALIGN.CENTER)
    footer(s, 16, TOTAL)

    # ── 17 Benefits ──
    s = prs.slides.add_slide(blank)
    bg(s, GREEN)
    title(s, "Why owners choose TableTap", "business impact")
    benefits = [
        ("Faster turnover", "Digital order + faster pay"),
        ("Fewer mistakes", "Tickets straight to kitchen"),
        ("One system", "Dine-in + takeaway + aggregators"),
        ("Less waste", "Stock + remote blocking"),
        ("Modern brand", "Premium guest UX"),
        ("Actionable data", "Daily CSV + team stats"),
    ]
    for i, (h, c) in enumerate(benefits):
        x = 0.8 + (i % 3) * 4.05
        y = 1.85 + (i // 3) * 1.75
        card(s, x, y, 3.55, 1.35, h, c, PANEL_2)
    stat_row(s, [("40%", "faster order", GREEN), ("0", "manual Swiggy entry", ORANGE), ("1", "platform", CYAN)], 5.35)
    footer(s, 17, TOTAL)

    # ── 18 CTA ──
    s = prs.slides.add_slide(blank)
    bg(s, ORANGE)
    textbox(s, "Ready for a live demo?", 0.8, 1.45, 11.7, 0.85, 46, WHITE, True, PP_ALIGN.CENTER)
    textbox(s, "Open a table · Scan the QR · Place an order · Watch the kitchen", 1.5, 2.45, 10.3, 0.55, 22, AMBER, True, PP_ALIGN.CENTER)
    card(s, 1.8, 3.65, 2.9, 1.15, "Staff login", "owner@varanasi.com", PANEL_2)
    card(s, 5.2, 3.65, 2.9, 1.15, "Password", "admin123", PANEL_2)
    card(s, 8.6, 3.65, 2.9, 1.15, "Super admin", "/platform/login", PANEL_2)
    card(s, 3.5, 5.15, 6.3, 1.0, "TableTap", "Scan · Order · Serve · Grow — the complete restaurant platform", PANEL_3)
    textbox(s, "tabletap · SETUP_GUIDE.md · AGGREGATOR_SETUP.md", 2.0, 6.55, 9.3, 0.35, 12, MUTED, False, PP_ALIGN.CENTER)
    footer(s, 18, TOTAL)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    prs.save(DETAILED_OUT)
    prs.save(LEGACY_OUT)
    print(f"Saved {OUT}")
    print(f"Saved {DETAILED_OUT}")
    print(f"Updated {LEGACY_OUT}")


if __name__ == "__main__":
    build()
